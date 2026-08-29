#!/usr/bin/env python3
"""Isolated offline worker for the Code-owned PPT Master pilot."""

from __future__ import annotations

import argparse
import csv
import hashlib
import importlib.util
import json
import os
import re
import sys
import types
import zipfile
from dataclasses import dataclass, field
from pathlib import Path


APP_ROOT = Path(__file__).resolve().parents[1]
if str(APP_ROOT) not in sys.path:
    sys.path.insert(0, str(APP_ROOT))

from scripts.validate_ppt_master_vendor import validate_vendor_package  # noqa: E402


SKILL_ROOT = APP_ROOT / "data" / "skills" / "ppt-master"
VENDOR_ROOT = SKILL_ROOT / "vendor"
FONT_FACE = "Microsoft YaHei"
INK = "18253B"
MUTED = "5F6B7A"
PAPER = "F6F3EC"
WHITE = "FFFFFF"
ORANGE = "FF6B35"
TEAL = "1E8A8A"
BLUE = "3165D4"


@dataclass
class Block:
    kind: str
    value: object


@dataclass
class SlideSpec:
    title: str
    blocks: list[Block] = field(default_factory=list)


def _audit_hook(event: str, args) -> None:
    if event.startswith("socket.") or event in {"subprocess.Popen", "os.system", "os.posix_spawn"}:
        raise RuntimeError(f"offline worker blocked audit event: {event}")


def _namespace(name: str, path: Path) -> None:
    module = types.ModuleType(name)
    module.__path__ = [str(path)]
    module.__package__ = name
    sys.modules[name] = module


def _load_module(name: str, path: Path):
    spec = importlib.util.spec_from_file_location(name, path)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"cannot load allowed vendor module: {path.name}")
    module = importlib.util.module_from_spec(spec)
    sys.modules[name] = module
    spec.loader.exec_module(module)
    return module


def _load_package(name: str, path: Path):
    spec = importlib.util.spec_from_file_location(
        name,
        path,
        submodule_search_locations=[str(path.parent)],
    )
    if spec is None or spec.loader is None:
        raise RuntimeError(f"cannot load allowed vendor package: {path.parent.name}")
    module = importlib.util.module_from_spec(spec)
    sys.modules[name] = module
    spec.loader.exec_module(module)
    return module


def _load_vendor_validators(manifest: dict):
    allowed = {item["path"] for item in manifest["files"] if item.get("category") == "python"}
    required = {
        "scripts/language_tags.py",
        "scripts/pptx_gradients.py",
        "scripts/pptx_shapes/__init__.py",
        "scripts/semantic_table.py",
        "scripts/svg_to_pptx/native_objects/chart_data.py",
        "scripts/svg_to_pptx/native_objects/table.py",
    }
    if not required <= allowed:
        raise RuntimeError("required vendor validator is outside the manifest")
    scripts = VENDOR_ROOT / "scripts"
    _namespace("code_ppt_master_vendor", scripts)
    _namespace("code_ppt_master_vendor.svg_to_pptx", scripts / "svg_to_pptx")
    _namespace(
        "code_ppt_master_vendor.svg_to_pptx.drawingml",
        scripts / "svg_to_pptx" / "drawingml",
    )
    _namespace(
        "code_ppt_master_vendor.svg_to_pptx.native_objects",
        scripts / "svg_to_pptx" / "native_objects",
    )
    _load_module("language_tags", scripts / "language_tags.py")
    _load_module("pptx_gradients", scripts / "pptx_gradients.py")
    _load_package("pptx_shapes", scripts / "pptx_shapes" / "__init__.py")
    _load_module("semantic_table", scripts / "semantic_table.py")
    chart_data = _load_module(
        "code_ppt_master_vendor.svg_to_pptx.native_objects.chart_data",
        scripts / "svg_to_pptx" / "native_objects" / "chart_data.py",
    )
    table = _load_module(
        "code_ppt_master_vendor.svg_to_pptx.native_objects.table",
        scripts / "svg_to_pptx" / "native_objects" / "table.py",
    )
    for module in list(sys.modules.values()):
        filename = getattr(module, "__file__", None)
        if not filename:
            continue
        path = Path(filename).resolve()
        try:
            relative = path.relative_to(VENDOR_ROOT).as_posix()
        except ValueError:
            continue
        if relative not in allowed:
            raise RuntimeError(f"worker imported an unmanifested vendor module: {relative}")
    return chart_data, table


def _clean_text(value: str, limit: int) -> str:
    text = re.sub(r"\s+", " ", value).strip()
    if not text or len(text) > limit:
        raise RuntimeError(f"slide text must be 1..{limit} characters")
    return text


def _pipe_cells(line: str) -> list[str]:
    return [cell.strip() for cell in line.strip().strip("|").split("|")]


def _table_at(lines: list[str], index: int):
    if index + 1 >= len(lines) or "|" not in lines[index]:
        return None
    separators = _pipe_cells(lines[index + 1])
    if not separators or not all(re.fullmatch(r":?-{3,}:?", item) for item in separators):
        return None
    headers = _pipe_cells(lines[index])
    if len(headers) != len(separators) or not 1 <= len(headers) <= 6:
        raise RuntimeError("markdown table must contain 1..6 columns")
    rows = []
    cursor = index + 2
    while cursor < len(lines) and "|" in lines[cursor] and lines[cursor].strip():
        row = _pipe_cells(lines[cursor])
        if len(row) != len(headers):
            raise RuntimeError("markdown table rows must match the header width")
        rows.append([_clean_text(cell, 80) for cell in row])
        cursor += 1
    if not 1 <= len(rows) <= 8:
        raise RuntimeError("markdown table must contain 1..8 data rows")
    return Block("table", {"headers": [_clean_text(cell, 80) for cell in headers], "rows": rows}), cursor


def _chart_block(lines: list[str], index: int):
    if lines[index].strip().lower() != "```chart":
        return None
    cursor, body = index + 1, []
    while cursor < len(lines) and lines[cursor].strip() != "```":
        body.append(lines[cursor])
        cursor += 1
    if cursor >= len(lines):
        raise RuntimeError("chart block is not closed")
    rows = list(csv.reader(body))
    if not 2 <= len(rows) <= 9 or not 2 <= len(rows[0]) <= 5:
        raise RuntimeError("chart block must contain a header plus 1..8 rows and 1..4 series")
    width = len(rows[0])
    if any(len(row) != width for row in rows):
        raise RuntimeError("chart CSV rows must have equal width")
    categories, series = [], [[] for _ in range(width - 1)]
    for row in rows[1:]:
        categories.append(_clean_text(row[0], 40))
        for offset, raw in enumerate(row[1:]):
            value = float(raw)
            if not (-1_000_000 <= value <= 1_000_000):
                raise RuntimeError("chart value is outside the supported range")
            series[offset].append(value)
    return Block("chart", {
        "categories": categories,
        "series": [
            {"name": _clean_text(rows[0][offset + 1], 40), "values": values}
            for offset, values in enumerate(series)
        ],
    }), cursor + 1


def parse_markdown(text: str) -> tuple[str, str, list[SlideSpec]]:
    lines = text.splitlines()
    title = "离线演示文稿"
    subtitle = "由 PPT Master 离线运行时生成"
    slides: list[SlideSpec] = []
    current: SlideSpec | None = None
    index = 0
    while index < len(lines):
        raw, stripped = lines[index], lines[index].strip()
        if stripped.startswith("# ") and not slides and current is None:
            title = _clean_text(stripped[2:], 80)
            index += 1
            continue
        if stripped.startswith("## "):
            if current:
                slides.append(current)
            current = SlideSpec(_clean_text(stripped[3:], 80))
            index += 1
            continue
        if current is None:
            if stripped:
                subtitle = _clean_text(stripped.lstrip("# "), 140)
            index += 1
            continue
        chart = _chart_block(lines, index)
        if chart:
            block, index = chart
            current.blocks.append(block)
            continue
        table = _table_at(lines, index)
        if table:
            block, index = table
            current.blocks.append(block)
            continue
        if stripped.startswith(("- ", "* ")):
            bullets = []
            while index < len(lines) and lines[index].strip().startswith(("- ", "* ")):
                bullets.append(_clean_text(lines[index].strip()[2:], 140))
                index += 1
            if len(bullets) > 6:
                raise RuntimeError("a slide supports at most 6 bullets")
            current.blocks.append(Block("bullets", bullets))
            continue
        if stripped:
            paragraph = _clean_text(stripped.lstrip("### "), 320)
            current.blocks.append(Block("paragraph", paragraph))
        index += 1
    if current:
        slides.append(current)
    if not slides:
        slides = [SlideSpec(title, [Block("paragraph", subtitle)])]
    if len(slides) > 11:
        raise RuntimeError("the offline pilot supports at most 11 content slides")
    for slide in slides:
        if len(slide.blocks) > 4:
            raise RuntimeError(f"slide is too dense: {slide.title}")
    return title, subtitle, slides


def _set_text_style(paragraph, size, color=INK, bold=False, align=None):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    paragraph.font.name = FONT_FACE
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = RGBColor.from_string(color)
    if align:
        paragraph.alignment = getattr(PP_ALIGN, align)


def _add_text(slide, text, x, y, w, h, *, size=20, color=INK, bold=False, align=None):
    from pptx.util import Inches

    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    _set_text_style(paragraph, size, color, bold, align)
    return shape


def _fill(shape, color, line=None, *, no_line=False):
    from pptx.dml.color import RGBColor

    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color)
    if no_line:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = RGBColor.from_string(line or color)


def _background(slide, slide_number: int):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.25), Inches(0.25), Inches(0.12), Inches(7.0)
    )
    _fill(bar, ORANGE, no_line=True)
    _add_text(slide, f"{slide_number:02d}", 12.35, 7.02, 0.55, 0.25, size=10, color=MUTED, align="RIGHT")


def _title_slide(prs, title: str, subtitle: str):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _background(slide, 1)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.9), Inches(0.9), Inches(2.3), Inches(2.3))
    _fill(circle, ORANGE)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(1.15), Inches(1.15), Inches(0.12))
    _fill(accent, TEAL)
    _add_text(slide, title, 0.9, 1.6, 8.8, 2.0, size=50, bold=True)
    _add_text(slide, subtitle, 0.95, 4.15, 7.7, 0.85, size=22, color=MUTED)
    _add_text(slide, "OFFLINE · EDITABLE · CJK", 0.95, 6.3, 4.4, 0.35, size=13, color=TEAL, bold=True)


def _add_table(slide, payload: dict, table_validator):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Inches, Pt

    table_validator._validate_table_payload({
        "schema": "ppt-master.semantic-table.v2",
        "columns": payload["headers"],
        "rows": payload["rows"],
        "header_rows": 1,
        "style": {},
    })
    rows, cols = len(payload["rows"]) + 1, len(payload["headers"])
    shape = slide.shapes.add_table(rows, cols, Inches(0.9), Inches(2.15), Inches(11.55), Inches(3.8))
    table = shape.table
    values = [payload["headers"], *payload["rows"]]
    for row_index, row in enumerate(values):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = value
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(
                INK if row_index == 0 else (WHITE if row_index % 2 else "E9EEF3")
            )
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.font.name = FONT_FACE
            paragraph.font.size = Pt(18 if row_index == 0 else 16)
            paragraph.font.bold = row_index == 0
            paragraph.font.color.rgb = RGBColor.from_string(WHITE if row_index == 0 else INK)


def _add_chart(slide, title: str, payload: dict, chart_validator):
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.util import Inches, Pt

    chart_payload = {
        "type": "column",
        "title": title,
        "categories": payload["categories"],
        "series": payload["series"],
    }
    chart_validator.validate_chart_payload(chart_payload)
    data = CategoryChartData()
    data.categories = payload["categories"]
    for series in payload["series"]:
        data.add_series(series["name"], series["values"])
    graphic = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.45), data,
    )
    chart = graphic.chart
    chart.has_legend = len(payload["series"]) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(13)
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.tick_labels.font.size = Pt(14)
    chart.category_axis.tick_labels.font.size = Pt(14)
    chart.chart_style = 10


def _content_slide(prs, spec: SlideSpec, number: int, validators):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _background(slide, number)
    marker = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(0.6), Inches(0.62), Inches(0.62))
    _fill(marker, TEAL)
    _add_text(slide, spec.title, 1.65, 0.52, 10.7, 0.75, size=36, bold=True)
    tables = [block for block in spec.blocks if block.kind == "table"]
    charts = [block for block in spec.blocks if block.kind == "chart"]
    if tables:
        _add_table(slide, tables[0].value, validators[1])
        return
    if charts:
        _add_chart(slide, spec.title, charts[0].value, validators[0])
        return
    bullets = []
    paragraphs = []
    for block in spec.blocks:
        if block.kind == "bullets":
            bullets.extend(block.value)
        elif block.kind == "paragraph":
            paragraphs.append(block.value)
    if paragraphs:
        _add_text(slide, "\n\n".join(paragraphs), 0.95, 1.75, 7.2, 3.7, size=22, color=INK)
    if bullets:
        from pptx.dml.color import RGBColor

        shape = _add_text(slide, bullets[0], 1.0, 1.7, 7.2, 4.8, size=22)
        frame = shape.text_frame
        frame.clear()
        from pptx.util import Pt
        for index, text in enumerate(bullets):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = text
            paragraph.level = 0
            paragraph.font.name = FONT_FACE
            paragraph.font.size = Pt(21)
            paragraph.font.color.rgb = RGBColor.from_string(INK)
            paragraph.space_after = Pt(13)
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(1.85), Inches(2.75), Inches(3.55))
    _fill(panel, INK)
    _add_text(slide, "结构化\n输入", 9.45, 2.35, 1.8, 0.95, size=26, color=WHITE, bold=True, align="CENTER")
    bridge = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(8.25), Inches(3.0), Inches(0.72), Inches(0.65))
    _fill(bridge, ORANGE)
    _add_text(slide, "原生对象\n输出", 9.45, 4.0, 1.8, 0.9, size=22, color="A7E2E2", bold=True, align="CENTER")


def _normalize_chart_axis_ids(pptx_path: Path) -> int:
    """Rewrite python-pptx's signed random axis ids to deterministic UInt32 values."""
    temporary = pptx_path.with_suffix(".axis-normalized.pptx")
    pattern = re.compile(rb'(<c:(?:axId|crossAx)\s+val=")(-?[0-9]+)(")')
    replacements = 0
    with zipfile.ZipFile(pptx_path) as source, zipfile.ZipFile(
        temporary, "w", compression=zipfile.ZIP_DEFLATED
    ) as target:
        for info in source.infolist():
            payload = source.read(info.filename)
            if re.fullmatch(r"ppt/charts/chart[0-9]+\.xml", info.filename):
                mapping: dict[bytes, bytes] = {}
                for match in pattern.finditer(payload):
                    raw = match.group(2)
                    if raw not in mapping:
                        mapping[raw] = str(1_000_000_000 + len(mapping) + 1).encode("ascii")
                payload, count = pattern.subn(
                    lambda match: match.group(1) + mapping[match.group(2)] + match.group(3),
                    payload,
                )
                replacements += count
            target.writestr(info, payload)
    os.replace(temporary, pptx_path)
    return replacements


def build_deck(request: dict, output: Path) -> dict:
    import importlib.metadata as metadata
    from pptx import Presentation
    from pptx.util import Inches

    expected = {item["project"]: item["version"] for item in request["expectedDependencies"]}
    for project, version in expected.items():
        if metadata.version(project) != version:
            raise RuntimeError(f"managed dependency version mismatch: {project}")
    vendor_summary = validate_vendor_package(SKILL_ROOT)
    manifest = json.loads((SKILL_ROOT / "vendor-manifest.json").read_text(encoding="utf-8"))
    if manifest.get("manifestDigest") != vendor_summary.get("manifestDigest"):
        raise RuntimeError("validated vendor manifest digest changed before import")
    chart_validator, table_validator = _load_vendor_validators(manifest)
    title, subtitle, slides = parse_markdown(request["markdown"])
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    _title_slide(prs, title, subtitle)
    for index, spec in enumerate(slides, start=2):
        _content_slide(prs, spec, index, (chart_validator, table_validator))
    temporary = output.with_suffix(".tmp.pptx")
    prs.save(temporary)
    normalized_axis_ids = _normalize_chart_axis_ids(temporary)
    os.replace(temporary, output)
    return {
        "schema": "code-ppt-master-worker-result/v1",
        "deckSha256": hashlib.sha256(output.read_bytes()).hexdigest(),
        "slideCount": len(prs.slides),
        "normalizedChartAxisIds": normalized_axis_ids,
        "vendorManifestDigest": vendor_summary["manifestDigest"],
        "vendorPrimitives": [
            "svg_to_pptx.native_objects.chart_data.validate_chart_payload",
            "svg_to_pptx.native_objects.table._validate_table_payload",
        ],
    }


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--request", required=True, type=Path)
    parser.add_argument("--output", required=True, type=Path)
    parser.add_argument("--result", required=True, type=Path)
    args = parser.parse_args()
    staging = args.request.resolve(strict=True).parent
    if args.output.resolve().parent != staging or args.result.resolve().parent != staging:
        raise RuntimeError("worker paths must share one run-owned staging directory")
    request = json.loads(args.request.read_text(encoding="utf-8"))
    if request.get("schema") != "code-ppt-master-worker/v1":
        raise RuntimeError("worker request schema is invalid")
    sys.addaudithook(_audit_hook)
    result = build_deck(request, args.output)
    args.result.write_text(
        json.dumps(result, ensure_ascii=False, sort_keys=True, indent=2) + "\n",
        encoding="utf-8",
        newline="\n",
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
