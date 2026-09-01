import base64
import hashlib
import io
import json
import os
import shutil
import sys
import tempfile
import threading
import unittest
import zipfile
from pathlib import Path
from unittest import mock

from code_runtime import ppt_master_runtime as runtime
import server as server_mod
from scripts import ppt_master_worker as worker


SAMPLE_MARKDOWN = """# 离线运行试点
结构化输入生成原生可编辑对象

## 安全边界保持明确
这一页先说明运行时的核心原则。
- 仅接受 UTF-8 Markdown 或 TXT
- 输出固定在当前 AgentRun
- 不访问网络、Key 或用户目录

## 原生对象保持可编辑
| 对象 | 验收状态 |
|---|---|
| 文本与形状 | 通过 |
| 表格与图表 | 通过 |

## 试点阶段稳步闭合
```chart
阶段,完成率
供应链,100
运行时,100
验证,100
```
"""


class TestPptMasterRuntime(unittest.TestCase):
    def setUp(self):
        self.temporary = tempfile.TemporaryDirectory()
        self.project = Path(self.temporary.name) / "project"
        self.project.mkdir()

    def tearDown(self):
        self.temporary.cleanup()

    def payload(self, *, run_id="ppt-run-1", operation="a" * 64, markdown=SAMPLE_MARKDOWN):
        return {
            "markdown": markdown,
            "_operationId": operation,
            "_projectRoot": str(self.project),
            "_agentRunId": run_id,
            "_toolCallId": "ppt-call-1",
            "_cancelEvent": threading.Event(),
        }

    @staticmethod
    def _record_digest(payload: bytes) -> str:
        encoded = base64.urlsafe_b64encode(hashlib.sha256(payload).digest()).rstrip(b"=")
        return "sha256=" + encoded.decode("ascii")

    def fake_dependency_contract(self):
        python_root = Path(self.temporary.name) / "managed-python"
        site_packages = python_root / "Lib" / "site-packages"
        site_packages.mkdir(parents=True)
        lock = {
            "schemaVersion": 1,
            "skill": "ppt-master",
            "capability": "offline-core",
            "wheels": [],
        }
        receipt_packages = []
        owned_files = {}
        specs = (
            ("skia-pathops", "0.9.2", "skia_pathops-0.9.2.dist-info", "pathops"),
            ("uharfbuzz", "0.50.0", "uharfbuzz-0.50.0.dist-info", "uharfbuzz"),
        )
        for project, version, dist_info, module in specs:
            module_dir = site_packages / module
            module_dir.mkdir()
            package_file = module_dir / "__init__.py"
            package_bytes = f"VERSION = {version!r}\n".encode("utf-8")
            package_file.write_bytes(package_bytes)
            dist_dir = site_packages / dist_info
            dist_dir.mkdir()
            metadata_file = dist_dir / "METADATA"
            metadata_bytes = f"Name: {project}\nVersion: {version}\n".encode("utf-8")
            metadata_file.write_bytes(metadata_bytes)
            record_file = dist_dir / "RECORD"
            rows = [
                [f"{module}/__init__.py", self._record_digest(package_bytes), str(len(package_bytes))],
                [f"{dist_info}/METADATA", self._record_digest(metadata_bytes), str(len(metadata_bytes))],
                [f"{dist_info}/RECORD", "", ""],
            ]
            record_text = "".join(
                ",".join(row) + "\n"
                for row in rows
            )
            record_file.write_text(record_text, encoding="utf-8", newline="")
            wheel_sha = hashlib.sha256(f"{project}=={version}".encode()).hexdigest()
            lock["wheels"].append({
                "project": project,
                "version": version,
                "sha256": wheel_sha,
            })
            receipt_packages.append({
                "project": project,
                "version": version,
                "wheelSha256": wheel_sha,
                "distInfo": dist_info,
                "recordSha256": hashlib.sha256(record_file.read_bytes()).hexdigest(),
            })
            owned_files[project] = package_file
        lock_path = python_root / "dependency-lock.json"
        lock_path.write_text(json.dumps(lock), encoding="utf-8", newline="\n")
        receipt = {
            "schema": "code-ppt-master-dependency-receipt/v1",
            "status": "installed",
            "skill": "ppt-master",
            "capability": "offline-core",
            "managedRuntime": "data/runtime/python",
            "lockSha256": hashlib.sha256(lock_path.read_bytes()).hexdigest(),
            "packages": receipt_packages,
        }
        receipt_path = python_root / "dependency-receipt.json"
        receipt_path.write_text(json.dumps(receipt), encoding="utf-8", newline="\n")
        return python_root, lock_path, receipt_path, owned_files

    def test_actual_worker_generates_editable_native_objects_and_replays_receipt(self):
        first = runtime.execute_ppt_master_tool(self.payload())
        self.assertTrue(first["ok"])
        self.assertFalse(first["replayed"])
        self.assertEqual(first["slideCount"], 4)
        self.assertGreater(first["shapeCount"], 10)
        self.assertEqual(first["tableCount"], 1)
        self.assertEqual(first["chartCount"], 1)
        self.assertTrue(first["editable"])
        self.assertTrue(first["offline"])
        deck = self.project / first["path"]
        self.assertTrue(deck.is_file())
        before = deck.read_bytes()
        with zipfile.ZipFile(deck) as archive:
            chart_xml = b"".join(
                archive.read(name)
                for name in archive.namelist()
                if name.startswith("ppt/charts/chart") and name.endswith(".xml")
            )
        self.assertNotRegex(chart_xml, rb'<c:(?:axId|crossAx) val="-')
        from pptx import Presentation

        presentation = Presentation(deck)
        layout_slide = presentation.slides[1]
        paragraph_shape = next(
            shape for shape in layout_slide.shapes
            if "这一页先说明" in getattr(shape, "text", "")
        )
        bullet_shape = next(
            shape for shape in layout_slide.shapes
            if "仅接受 UTF-8" in getattr(shape, "text", "")
        )
        self.assertLessEqual(
            paragraph_shape.top + paragraph_shape.height,
            bullet_shape.top,
        )

        replay = runtime.execute_ppt_master_tool(self.payload())
        self.assertTrue(replay["ok"])
        self.assertTrue(replay["replayed"])
        self.assertEqual(replay["sha256"], first["sha256"])
        self.assertEqual(deck.read_bytes(), before)
        self.assertFalse((self.project / "output/ppt-master/.staging/ppt-run-1").exists())
        self.assertFalse((self.project / "%SystemDrive%").exists())

    def test_dependency_record_tamper_hides_tool_and_starts_no_worker(self):
        python_root, lock_path, receipt_path, files = self.fake_dependency_contract()
        patches = mock.patch.multiple(
            runtime,
            PYTHON_ROOT=python_root,
            LOCK_PATH=lock_path,
            DEPENDENCY_RECEIPT_PATH=receipt_path,
            EXPECTED_DEPENDENCY_RECEIPT_DIGEST=hashlib.sha256(
                json.dumps(
                    json.loads(receipt_path.read_text(encoding="utf-8")),
                    ensure_ascii=False,
                    sort_keys=True,
                    separators=(",", ":"),
                ).encode("utf-8")
            ).hexdigest(),
        )
        with patches:
            self.assertTrue(runtime.validate_ppt_master_dependency_installation()["ok"])
            files["skia-pathops"].write_text("tampered\n", encoding="utf-8")
            with mock.patch.object(
                server_mod,
                "inspect_skill_directory",
                return_value={"status": "ready"},
            ):
                skills = {item["name"]: item for item in server_mod.list_skills(brief=True)}
                capability = server_mod.get_single_skill_dependency_status(
                    "ppt-master", "offline-core"
                )
            self.assertEqual(skills["ppt-master"]["tools"], [])
            self.assertEqual(capability["status"], "unavailable")
            with mock.patch.object(runtime.subprocess, "Popen") as popen:
                with self.assertRaises(runtime.PptMasterRuntimeError) as raised:
                    runtime.execute_ppt_master_tool(self.payload(run_id="tampered-runtime"))
            self.assertEqual(raised.exception.code, "ppt_master_dependency_integrity")
            popen.assert_not_called()

    def test_dependency_record_and_extra_target_distribution_fail_closed(self):
        python_root, lock_path, receipt_path, _ = self.fake_dependency_contract()
        with mock.patch.multiple(
            runtime,
            PYTHON_ROOT=python_root,
            LOCK_PATH=lock_path,
            DEPENDENCY_RECEIPT_PATH=receipt_path,
            EXPECTED_DEPENDENCY_RECEIPT_DIGEST=hashlib.sha256(
                json.dumps(
                    json.loads(receipt_path.read_text(encoding="utf-8")),
                    ensure_ascii=False,
                    sort_keys=True,
                    separators=(",", ":"),
                ).encode("utf-8")
            ).hexdigest(),
        ):
            record = (
                python_root / "Lib/site-packages/skia_pathops-0.9.2.dist-info/RECORD"
            )
            original = record.read_bytes()
            record.write_bytes(original + b"\n")
            with self.assertRaises(runtime.PptMasterRuntimeError) as raised:
                runtime.validate_ppt_master_dependency_installation()
            self.assertEqual(raised.exception.code, "ppt_master_dependency_integrity")

            record.write_bytes(original)
            source = record.parent
            duplicate = source.parent / "skia_pathops-9.9.9.dist-info"
            shutil.copytree(source, duplicate)
            with self.assertRaises(runtime.PptMasterRuntimeError) as raised:
                runtime.validate_ppt_master_dependency_installation()
            self.assertEqual(raised.exception.code, "ppt_master_dependency_integrity")

    def test_final_receipt_and_deck_reparse_are_rejected_before_read(self):
        from pptx import Presentation

        original_is_reparse = runtime._is_reparse
        original_read_text = Path.read_text
        original_open = Path.open
        for filename in ("receipt.json", "presentation.pptx"):
            with self.subTest(filename=filename):
                final_dir = self.project / "output" / "ppt-master" / f"final-{filename.split('.')[0]}"
                final_dir.mkdir(parents=True)
                deck = final_dir / "presentation.pptx"
                presentation = Presentation()
                presentation.slides.add_slide(presentation.slide_layouts[6])
                presentation.save(deck)
                operation = "e" * 64
                markdown_sha = hashlib.sha256(b"# final").hexdigest()
                (final_dir / "receipt.json").write_text(json.dumps({
                    "schema": "code-ppt-master-output-receipt/v1",
                    "operationId": operation,
                    "inputSha256": markdown_sha,
                    "deckSha256": hashlib.sha256(deck.read_bytes()).hexdigest(),
                }), encoding="utf-8")
                marked = final_dir / filename

                def guarded_read_text(path, *args, **kwargs):
                    if path == marked:
                        raise AssertionError("reparse receipt was read")
                    return original_read_text(path, *args, **kwargs)

                def guarded_open(path, *args, **kwargs):
                    if path == marked:
                        raise AssertionError("reparse deck was opened")
                    return original_open(path, *args, **kwargs)

                with mock.patch.object(
                    runtime,
                    "_is_reparse",
                    side_effect=lambda path: path == marked or original_is_reparse(path),
                ), mock.patch.object(Path, "read_text", guarded_read_text), mock.patch.object(
                    Path, "open", guarded_open
                ):
                    with self.assertRaises(runtime.PptMasterRuntimeError) as raised:
                        runtime._recover_final(final_dir, operation, markdown_sha, "presentation.pptx")
                self.assertEqual(raised.exception.code, "ppt_master_reparse_blocked")

    def test_existing_staging_files_reparse_are_rejected_before_read_or_replay(self):
        from pptx import Presentation

        original_is_reparse = runtime._is_reparse
        original_read_text = Path.read_text
        original_open = Path.open
        for index, filename in enumerate((
            "prepared.json", "request.json", "worker-result.json", "presentation.pptx"
        )):
            with self.subTest(filename=filename):
                run_id = f"staging-reparse-{index}"
                operation = f"{index + 1:x}" * 64
                input_sha = hashlib.sha256(SAMPLE_MARKDOWN.encode("utf-8")).hexdigest()
                staging = self.project / "output" / "ppt-master" / ".staging" / run_id
                staging.mkdir(parents=True)
                deck = staging / "presentation.pptx"
                presentation = Presentation()
                presentation.slides.add_slide(presentation.slide_layouts[6])
                presentation.save(deck)
                (staging / "prepared.json").write_text(json.dumps({
                    "operationId": operation,
                    "inputSha256": input_sha,
                }), encoding="utf-8")
                (staging / "request.json").write_text("{}", encoding="utf-8")
                (staging / "worker-result.json").write_text(json.dumps({
                    "deckSha256": hashlib.sha256(deck.read_bytes()).hexdigest(),
                }), encoding="utf-8")
                marked = staging / filename

                def guarded_read_text(path, *args, **kwargs):
                    if path == marked:
                        raise AssertionError("reparse staging file was read")
                    return original_read_text(path, *args, **kwargs)

                def guarded_open(path, *args, **kwargs):
                    if path == marked:
                        raise AssertionError("reparse staging file was opened")
                    return original_open(path, *args, **kwargs)

                with mock.patch.object(
                    runtime, "_validate_runtime_contract", return_value={
                        "receipt": {"packages": []},
                        "vendor": {"manifestDigest": "x"},
                    }
                ), mock.patch.object(
                    runtime,
                    "_is_reparse",
                    side_effect=lambda path: path == marked or original_is_reparse(path),
                ), mock.patch.object(Path, "read_text", guarded_read_text), mock.patch.object(
                    Path, "open", guarded_open
                ), mock.patch.object(runtime.subprocess, "Popen") as popen:
                    with self.assertRaises(runtime.PptMasterRuntimeError) as raised:
                        runtime.execute_ppt_master_tool(self.payload(
                            run_id=run_id,
                            operation=operation,
                        ))
                self.assertEqual(raised.exception.code, "ppt_master_reparse_blocked")
                popen.assert_not_called()

    def test_markdown_paragraphs_merge_and_structured_mixing_fails_closed(self):
        _, _, slides = worker.parse_markdown(
            "# Deck\nSubtitle\n\n## Layout\nFirst ordinary line\nSecond ordinary line\n- bullet one\n- bullet two\n"
        )
        self.assertEqual([block.kind for block in slides[0].blocks], ["paragraph", "bullets"])
        self.assertEqual(slides[0].blocks[0].value, "First ordinary line Second ordinary line")
        for structured in (
            "| A | B |\n|---|---|\n| 1 | 2 |",
            "```chart\nStage,Value\nDone,100\n```",
        ):
            with self.subTest(structured=structured), self.assertRaisesRegex(
                RuntimeError, "cannot mix"
            ):
                worker.parse_markdown(
                    "# Deck\nSubtitle\n\n## Mixed\nThis paragraph must not disappear.\n" + structured
                )

    def test_ooxml_validator_rejects_forbidden_internal_parts_and_content_types(self):
        from pptx import Presentation
        from xml.etree import ElementTree as ET

        base = Path(self.temporary.name) / "base.pptx"
        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[6])
        presentation.save(base)
        cases = (
            ("ppt/vbaProject.bin", "application/vnd.ms-office.vbaProject"),
            ("ppt/embeddings/evil.bin", "application/vnd.openxmlformats-officedocument.oleObject"),
            ("ppt/activeX/activeX1.bin", "application/vnd.ms-office.activeX"),
            ("ppt/media/audio1.mp3", "audio/mpeg"),
            ("ppt/media/video1.mp4", "video/mp4"),
            ("customXml/item1.xml", "application/xml"),
        )
        namespace = "http://schemas.openxmlformats.org/package/2006/content-types"
        for index, (part_name, content_type) in enumerate(cases):
            with self.subTest(part_name=part_name):
                candidate = Path(self.temporary.name) / f"forbidden-{index}.pptx"
                with zipfile.ZipFile(base) as source, zipfile.ZipFile(candidate, "w") as target:
                    for info in source.infolist():
                        payload = source.read(info.filename)
                        if info.filename == "[Content_Types].xml":
                            root = ET.fromstring(payload)
                            ET.SubElement(root, f"{{{namespace}}}Override", {
                                "PartName": "/" + part_name,
                                "ContentType": content_type,
                            })
                            payload = ET.tostring(root, encoding="utf-8", xml_declaration=True)
                        target.writestr(info, payload)
                    target.writestr(part_name, b"blocked")
                with self.assertRaises(runtime.PptMasterRuntimeError) as raised:
                    runtime._validate_pptx(candidate)
                self.assertEqual(raised.exception.code, "ppt_master_forbidden_part")

    def test_ooxml_validator_rejects_macro_inside_allowed_native_chart_workbook(self):
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches
        from xml.etree import ElementTree as ET

        source_path = Path(self.temporary.name) / "chart-source.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        chart_data = CategoryChartData()
        chart_data.categories = ["A", "B"]
        chart_data.add_series("Value", [1, 2])
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1), Inches(1), Inches(8), Inches(4), chart_data,
        )
        presentation.save(source_path)
        candidate = Path(self.temporary.name) / "chart-with-macro.pptx"
        namespace = "http://schemas.openxmlformats.org/package/2006/content-types"
        with zipfile.ZipFile(source_path) as source, zipfile.ZipFile(candidate, "w") as target:
            for info in source.infolist():
                payload = source.read(info.filename)
                if info.filename.startswith("ppt/embeddings/") and info.filename.endswith(".xlsx"):
                    output = io.BytesIO()
                    with zipfile.ZipFile(io.BytesIO(payload)) as workbook, zipfile.ZipFile(output, "w") as rewritten:
                        for workbook_info in workbook.infolist():
                            workbook_payload = workbook.read(workbook_info.filename)
                            if workbook_info.filename == "[Content_Types].xml":
                                root = ET.fromstring(workbook_payload)
                                ET.SubElement(root, f"{{{namespace}}}Override", {
                                    "PartName": "/xl/vbaProject.bin",
                                    "ContentType": "application/vnd.ms-office.vbaProject",
                                })
                                workbook_payload = ET.tostring(
                                    root, encoding="utf-8", xml_declaration=True
                                )
                            rewritten.writestr(workbook_info, workbook_payload)
                        rewritten.writestr("xl/vbaProject.bin", b"blocked")
                    payload = output.getvalue()
                target.writestr(info, payload)
        with self.assertRaises(runtime.PptMasterRuntimeError) as raised:
            runtime._validate_pptx(candidate)
        self.assertEqual(raised.exception.code, "ppt_master_forbidden_part")

    def test_inline_and_project_file_inputs_fail_closed(self):
        normalized = runtime._normalize_markdown({"markdown": "# 标题"}, self.project)
        self.assertEqual(normalized[:2], ("# 标题", "inline"))
        source = self.project / "brief.md"
        source.write_text("# 项目简报", encoding="utf-8")
        normalized = runtime._normalize_markdown({"sourcePath": "brief.md"}, self.project)
        self.assertEqual(normalized[:2], ("# 项目简报", "project_file"))

        invalid = (
            {},
            {"markdown": "# A", "sourcePath": "brief.md"},
            {"markdown": "see https://example.com"},
            {"markdown": '<img src="https://example.com/a.png">'},
            {"sourcePath": str(source)},
            {"sourcePath": "../brief.md"},
            {"sourcePath": "brief.pdf"},
        )
        for payload in invalid:
            with self.subTest(payload=payload), self.assertRaises(runtime.PptMasterRuntimeError):
                runtime._normalize_markdown(payload, self.project)

        source.write_bytes(b"x" * (runtime.MAX_INPUT_BYTES + 1))
        with self.assertRaises(runtime.PptMasterRuntimeError) as raised:
            runtime._normalize_markdown({"sourcePath": "brief.md"}, self.project)
        self.assertEqual(raised.exception.code, "ppt_master_source_too_large")

    def test_reparse_source_is_blocked_without_reading_it(self):
        source = self.project / "brief.txt"
        source.write_text("protected", encoding="utf-8")
        original = runtime._is_reparse
        with mock.patch.object(
            runtime,
            "_is_reparse",
            side_effect=lambda path: path == source or original(path),
        ):
            with self.assertRaises(runtime.PptMasterRuntimeError) as raised:
                runtime._normalize_markdown({"sourcePath": "brief.txt"}, self.project)
        self.assertEqual(raised.exception.code, "ppt_master_reparse_blocked")

    def test_existing_public_output_is_never_overwritten(self):
        occupied = self.project / "output/ppt-master/ppt-run-occupied"
        occupied.mkdir(parents=True)
        sentinel = occupied / "sentinel.bin"
        sentinel.write_bytes(b"keep")
        with self.assertRaises(runtime.PptMasterRuntimeError) as raised:
            runtime.execute_ppt_master_tool(self.payload(run_id="ppt-run-occupied"))
        self.assertEqual(raised.exception.code, "ppt_master_output_exists")
        self.assertEqual(sentinel.read_bytes(), b"keep")

    def test_restart_with_prepared_only_is_not_replayed_or_regenerated(self):
        run_id, operation = "ppt-run-unknown", "b" * 64
        staging = self.project / f"output/ppt-master/.staging/{run_id}"
        staging.mkdir(parents=True)
        input_sha = hashlib.sha256(SAMPLE_MARKDOWN.encode("utf-8")).hexdigest()
        (staging / "prepared.json").write_text(
            json.dumps({"operationId": operation, "inputSha256": input_sha}),
            encoding="utf-8",
        )
        with mock.patch.object(runtime, "_validate_runtime_contract", return_value={
            "receipt": {"packages": []},
            "vendor": {"manifestDigest": "x"},
        }), mock.patch.object(runtime.subprocess, "Popen") as popen:
            with self.assertRaises(runtime.PptMasterRuntimeError) as raised:
                runtime.execute_ppt_master_tool(self.payload(run_id=run_id, operation=operation))
        self.assertEqual(raised.exception.code, "ppt_master_previous_outcome_unknown")
        self.assertTrue(raised.exception.outcome_unknown)
        popen.assert_not_called()
        self.assertFalse(staging.exists())

    def test_cancel_and_timeout_kill_worker_and_remove_staging(self):
        fake = Path(self.temporary.name) / "slow_worker.py"
        fake.write_text("import time\ntime.sleep(60)\n", encoding="utf-8")
        contract = {
            "receipt": {"packages": []},
            "vendor": {"manifestDigest": "x"},
        }
        for mode in ("cancel", "timeout"):
            run_id = f"ppt-{mode}"
            payload = self.payload(run_id=run_id, operation=("c" if mode == "cancel" else "d") * 64)
            if mode == "cancel":
                payload["_cancelEvent"].set()
            with self.subTest(mode=mode), mock.patch.object(
                runtime, "_validate_runtime_contract", return_value=contract
            ), mock.patch.object(runtime, "MANAGED_PYTHON", Path(sys.executable)), mock.patch.object(
                runtime, "WORKER_PATH", fake
            ), mock.patch.object(runtime, "TIMEOUT_SECONDS", 0.05):
                with self.assertRaises(runtime.PptMasterRuntimeError) as raised:
                    runtime.execute_ppt_master_tool(payload)
            self.assertEqual(raised.exception.cancelled, mode == "cancel")
            self.assertEqual(raised.exception.timed_out, mode == "timeout")
            self.assertFalse((self.project / f"output/ppt-master/{run_id}").exists())
            self.assertFalse((self.project / f"output/ppt-master/.staging/{run_id}").exists())

    def test_tool_schema_exposes_no_command_module_or_output_path(self):
        definition = server_mod._SERVER_TOOL_DEFINITIONS["create_ppt_master_deck"]
        properties = definition["function"]["parameters"]["properties"]
        self.assertEqual(set(properties), {"markdown", "sourcePath"})
        errors = server_mod._registered_tool_argument_errors(
            "create_ppt_master_deck",
            {"markdown": "# A", "outputPath": "outside.pptx", "module": "anything"},
        )
        self.assertTrue(errors)
        payload = {"tools": [definition]}
        for profile in ("read", "plan"):
            self.assertEqual(
                server_mod._agent_selected_tools(payload, ["create_ppt_master_deck"], profile),
                [],
            )
        for profile in ("accept", "bypass"):
            selected = server_mod._agent_selected_tools(
                payload, ["create_ppt_master_deck"], profile,
            )
            self.assertEqual(selected[0]["function"]["name"], "create_ppt_master_deck")


if __name__ == "__main__":
    unittest.main()
