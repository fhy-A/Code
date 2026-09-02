import hashlib
import importlib.util
import json
import shutil
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path
from unittest import mock

import server as server_mod
from code_runtime.skill_dependencies import inspect_skill_directory
from code_runtime.skill_resources import (
    SkillResourceError,
    audit_bundled_skill_resources,
    resolve_skill_resources,
)


ROOT = Path(__file__).resolve().parents[1]
PPTX_DIR = ROOT / "data" / "skills" / "pptx"
VALIDATE_PATH = PPTX_DIR / "scripts" / "office" / "validate.py"
RENDER_PATH = PPTX_DIR / "scripts" / "render.py"
MANIFEST_PATH = PPTX_DIR / "code-resources.json"


def _sha256(path):
    return hashlib.sha256(Path(path).read_bytes().replace(b"\r\n", b"\n")).hexdigest()


def _load_module(name, path):
    spec = importlib.util.spec_from_file_location(name, path)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


class TestPptxSkillResourceContract(unittest.TestCase):
    def test_repository_contract_projects_exact_resources_and_strict_render_dependency(self):
        manifest = json.loads(MANIFEST_PATH.read_text(encoding="utf-8"))
        resources = {item["id"]: item for item in manifest["resources"]}
        self.assertEqual(manifest["schemaVersion"], 1)
        self.assertEqual(manifest["skill"], "pptx")
        self.assertEqual(resources["validate-deck"]["sha256"], _sha256(VALIDATE_PATH))
        self.assertEqual(resources["render-pdf"]["sha256"], _sha256(RENDER_PATH))

        resolved = resolve_skill_resources("pptx", ROOT / "data" / "skills", ROOT / "data" / "skills")
        self.assertEqual(resolved["source"], "installed")
        self.assertEqual(
            {item["id"] for item in resolved["resources"]},
            {"validate-deck", "render-pdf"},
        )
        self.assertTrue(all(Path(item["path"]).is_relative_to(PPTX_DIR) for item in resolved["resources"]))
        self.assertIn("do not search", resolved["instructions"].lower())

        dependencies = json.loads((PPTX_DIR / "dependencies.json").read_text(encoding="utf-8"))
        render_soffice = next(
            item for item in dependencies["capabilities"]["render"]["required"]
            if item["name"] == "soffice"
        )
        self.assertEqual(render_soffice["executableKind"], "libreoffice")
        skill_text = (PPTX_DIR / "SKILL.md").read_text(encoding="utf-8")
        self.assertIn("runtimeResources", skill_text)
        self.assertNotIn("claude-skills", skill_text)
        self.assertNotIn("scripts/", skill_text)
        audit = audit_bundled_skill_resources(ROOT / "data" / "skills")
        self.assertTrue(audit["ok"])
        self.assertEqual(audit["findings"], [])

    def test_same_root_code_owned_pptx_and_xlsx_contracts_keep_integrity_gates(self):
        with tempfile.TemporaryDirectory(prefix="Code same-root bundled resource integrity ") as temp:
            shared_skills = Path(temp) / "Code Dev with spaces" / "data" / "skills"
            for name in ("pptx", "xlsx"):
                source_skill = ROOT / "data" / "skills" / name
                copied_skill = shared_skills / name
                shutil.copytree(
                    source_skill,
                    copied_skill,
                    ignore=shutil.ignore_patterns("__pycache__", "*.pyc"),
                )
                resolved = resolve_skill_resources(name, shared_skills, shared_skills)
                self.assertEqual(resolved["source"], "installed")

                manifest_path = copied_skill / "code-resources.json"
                manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
                relative = manifest["resources"][0]["path"]
                helper = copied_skill / relative
                helper.unlink()
                with self.assertRaises(SkillResourceError) as missing:
                    resolve_skill_resources(name, shared_skills, shared_skills)
                self.assertEqual(missing.exception.error_code, "bundled_resource_missing")

                shutil.copy2(source_skill / relative, helper)
                manifest.pop("compatibleInstalled")
                manifest_path.write_text(json.dumps(manifest), encoding="utf-8")
                with self.assertRaises(SkillResourceError) as malformed:
                    resolve_skill_resources(name, shared_skills, shared_skills)
                self.assertEqual(malformed.exception.error_code, "resource_contract_invalid")

    def test_legacy_installed_copy_uses_read_only_bundled_fallback_and_server_projects_it(self):
        with tempfile.TemporaryDirectory(prefix="Code PPTX app-style resources ") as temp:
            root = Path(temp)
            bundle_root = root / "app bundle" / "data" / "skills"
            installed_root = root / "moved active skills" / "skills"
            shutil.copytree(PPTX_DIR, bundle_root / "pptx")
            installed_skill = installed_root / "pptx"
            installed_skill.mkdir(parents=True)
            for name in ("SKILL.md", "dependencies.json"):
                shutil.copy2(PPTX_DIR / name, installed_skill / name)
            before = {
                path.relative_to(installed_skill).as_posix(): path.read_bytes()
                for path in installed_skill.rglob("*")
                if path.is_file()
            }

            resolved = resolve_skill_resources("pptx", installed_root, bundle_root)
            with (
                mock.patch.object(server_mod, "APP_DIR", root / "app bundle"),
                mock.patch.object(server_mod, "SKILLS_DIR", installed_root),
            ):
                projected = server_mod.execute_use_skill_tool({"name": "pptx"})

            after = {
                path.relative_to(installed_skill).as_posix(): path.read_bytes()
                for path in installed_skill.rglob("*")
                if path.is_file()
            }
        self.assertEqual(resolved["source"], "bundled-fallback")
        self.assertTrue(all(Path(item["path"]).is_relative_to(bundle_root / "pptx") for item in resolved["resources"]))
        self.assertEqual(before, after)
        self.assertTrue(projected["ok"])
        self.assertEqual(projected["runtimeResources"]["source"], "bundled-fallback")
        self.assertIn("do not search", projected["runtimeResources"]["instructions"].lower())


class TestPptxCleanRoomResources(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        scripts_dir = str(RENDER_PATH.parent)
        if scripts_dir not in sys.path:
            sys.path.insert(0, scripts_dir)
        cls.validator = _load_module("code_pptx_validate", VALIDATE_PATH)
        cls.renderer = _load_module("code_pptx_render", RENDER_PATH)

    def test_controlled_pptx_creation_validates_renders_and_has_a_readable_page(self):
        from PIL import Image
        from PyPDF2 import PdfReader
        from pptx import Presentation

        with tempfile.TemporaryDirectory(prefix="Code PPTX clean room ") as temp:
            root = Path(temp)
            deck = root / "controlled deck.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[0])
            slide.shapes.title.text = "Northstar Studio"
            slide.placeholders[1].text = "Controlled PPTX resource validation"
            presentation.save(deck)

            validation = self.validator.validate_presentation(deck)
            inspection = inspect_skill_directory(
                PPTX_DIR,
                bundled_skills_dir=ROOT / "data" / "skills",
                app_dir=ROOT,
                data_dir=ROOT / "data",
                capability_id="render",
            )
            selected = next(item for item in inspection["capabilities"] if item["id"] == "render")
            executables = {item["name"]: item["executable"] for item in selected["required"]}
            rendered = self.renderer.render_presentation(
                deck,
                root / "rendered output",
                timeout_seconds=60,
                soffice_path=executables["soffice"],
            )
            pdf = Path(rendered["pdf"])
            reader = PdfReader(str(pdf))
            prefix = root / "rendered output" / "slide"
            subprocess.run(
                [executables["pdftoppm"], "-jpeg", "-r", "72", str(pdf), str(prefix)],
                check=True,
                capture_output=True,
                timeout=30,
            )
            images = list((root / "rendered output").glob("slide-*.jpg"))
            self.assertEqual(validation["status"], "success")
            self.assertEqual(validation["slides"], 1)
            self.assertEqual(rendered["status"], "success")
            self.assertEqual(len(reader.pages), 1)
            self.assertEqual(len(images), 1)
            with Image.open(images[0]) as image:
                self.assertGreater(image.width, 100)
                self.assertGreater(image.height, 100)


if __name__ == "__main__":
    unittest.main()
